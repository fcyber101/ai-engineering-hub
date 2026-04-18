

from loguru import logger
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
from typing import Dict, Any
import io

from core.state import PresentationState


def add_glass_shape(slide, left, top, width, height, transparency=60):
    """Add glass shape"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )

    try:
        from pptx.oxml.ns import qn
        from pptx.oxml.xmlchemy import OxmlElement

        fill = shape.fill
        fill.solid()
        spPr = shape._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = OxmlElement('a:solidFill')
            spPr.append(solidFill)
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is None:
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', 'FFFFFF')
            solidFill.append(srgbClr)
        alpha = OxmlElement('a:alpha')
        alpha_val = int((transparency / 100) * 100000)
        alpha.set('val', str(alpha_val))
        existing_alpha = srgbClr.find(qn('a:alpha'))
        if existing_alpha is not None:
            srgbClr.remove(existing_alpha)
        srgbClr.append(alpha)
    except Exception:
        pass
    return shape

def node_6_2_ppt_generation(state: PresentationState) -> PresentationState:
    """Generate PPTX using images from memory"""

    logger.info("Node 6.2: Creating presentation from in-memory images")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Get image bytes from state
    image_bytes = state.get("image_bytes", {})
    formatted_slides = state.get("formatted_slides", [])
    structured_sections = state.get("structured_sections", [])
    theme = state.get("presentation_theme", "Presentation")

    # Find blank layout
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    # Title Slide
    logger.info("Creating Title Slide...")
    slide = prs.slides.add_slide(blank_layout)

    # Add background from memory bytes
    if "title_background" in image_bytes:
        img_stream = io.BytesIO(image_bytes["title_background"])
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)
    else:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 39, 97)

    # Title text
    title_text = formatted_slides[0]["title"] if formatted_slides else theme
    add_glass_shape(slide, Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5), transparency=60)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(30, 39, 97)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Content Slide
    logger.info(f"Creating {len(formatted_slides)} content slides...")
    content_bytes = image_bytes.get("content_background")

    for idx, slide_data in enumerate(formatted_slides):
        slide = prs.slides.add_slide(blank_layout)

        # Add background from memory
        if content_bytes:
            img_stream = io.BytesIO(content_bytes)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0),
                                     width=prs.slide_width, height=prs.slide_height)
        else:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 245, 245)

        # Title bar
        title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                          Inches(13.333), Inches(0.9))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(30, 39, 97)
        title_bg.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get("title", f"Slide {idx + 1}")
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Bullets and summary
        bullets = slide_data.get("bullets", [])
        section_summary = ""
        if idx < len(structured_sections):
            section_summary = structured_sections[idx].get("section_summary", "")

        if bullets or section_summary:
            content_glass = add_glass_shape(slide, Inches(0.3), Inches(1.05), Inches(12.733), Inches(5.7), transparency=50)
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.2), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for bullet in bullets:
                p = content_frame.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.space_after = Pt(10)

            if bullets and section_summary:
                spacer = content_frame.add_paragraph()
                spacer.text = ""
                spacer.space_after = Pt(15)

            if section_summary:
                p_sep = content_frame.add_paragraph()
                p_sep.text = "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"
                p_sep.font.size = Pt(12)
                p_sep.font.color.rgb = RGBColor(150, 150, 150)
                p_sep.space_after = Pt(12)

                p_summary = content_frame.add_paragraph()
                p_summary.text = section_summary
                p_summary.font.size = Pt(12)
                p_summary.font.italic = True
                p_summary.font.color.rgb = RGBColor(60, 60, 60)

    # Summary Slide
    logger.info("Creating Summary Slide...")
    slide = prs.slides.add_slide(blank_layout)

    if content_bytes:
        img_stream = io.BytesIO(content_bytes)
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(13.333), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(30, 39, 97)
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Key Takeaways"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_glass_shape(slide, Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5), transparency=50)

    executive_summary = state.get("executive_summary", "Key insights from this presentation.")
    summary_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.2))
    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True

    p = summary_frame.add_paragraph()
    p.text = executive_summary
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(30, 39, 97)
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(20)

    # Thank you Slide
    logger.info("Creating Thank You Slide...")
    slide = prs.slides.add_slide(blank_layout)

    if "thanks_background" in image_bytes:
        img_stream = io.BytesIO(image_bytes["thanks_background"])
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)
    else:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 39, 97)

    add_glass_shape(slide, Inches(1.5), Inches(2.2), Inches(10.333), Inches(2.8), transparency=50)

    thanks_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.333), Inches(2.8))
    thanks_frame = thanks_box.text_frame
    thanks_frame.word_wrap = True
    thanks_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = thanks_frame.add_paragraph()
    p1.text = "Thank You!"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(30, 39, 97)
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(20)

    p2 = thanks_frame.add_paragraph()
    p2.text = "Questions?"
    p2.font.size = Pt(32)
    p2.font.bold = False
    p2.font.color.rgb = RGBColor(30, 39, 97)
    p2.alignment = PP_ALIGN.CENTER

    # Save final presentation
    Path("outputs").mkdir(exist_ok=True)
    ppt_path = "outputs/generated_presentation.pptx"
    prs.save(ppt_path)
    logger.success(f"PPTX saved: {ppt_path}")

    return {
        **state,
        "ppt_file_path": ppt_path
    }